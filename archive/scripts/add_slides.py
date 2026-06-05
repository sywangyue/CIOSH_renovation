#!/usr/bin/env python3
"""Add S07 and S08 slides to CIOSH Blackbox PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

CIOSH_GREEN = RGBColor(0x00, 0x97, 0x48)
CIOSH_ORANGE = RGBColor(0xF3, 0x97, 0x00)
DARK = RGBColor(0x1A, 0x1A, 0x1A)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation("/Volumes/databoard/AI Project/ciosh/output/CIOSH_Blackbox_拆解路径.pptx")
blank_layout = prs.slide_layouts[7]  # 空白

def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multiline(slide, left, top, width, height, lines, font_size=12, color=DARK, line_spacing=1.3):
    """Add multiline text box. lines is list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text = line_data[0]
        bold = line_data[1] if len(line_data) > 1 else False
        clr = line_data[2] if len(line_data) > 2 else color
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = clr
        p.font.name = 'Arial'
        p.space_after = Pt(4)
    return tf

def add_rect(slide, left, top, width, height, fill_color, text='', font_size=11, font_color=WHITE):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    return shape

def add_big_number(slide, left, top, number, label, num_color=CIOSH_GREEN):
    add_textbox(slide, left, top, 2.5, 0.6, str(number), font_size=36, bold=True, color=num_color, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top+0.55, 2.5, 0.4, label, font_size=10, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: S07 封面 / 核心数据
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)

# Top accent bar
add_rect(slide2, 0, 0, 13.33, 0.06, CIOSH_GREEN)
# Side accent bar
add_rect(slide2, 0, 0.06, 0.06, 7.44, CIOSH_ORANGE)

# Tag
add_textbox(slide2, 0.5, 0.25, 2, 0.3, "S07 · 需求验证", font_size=10, bold=True, color=CIOSH_GREEN)
# Title
add_textbox(slide2, 0.5, 0.55, 11, 0.8, "劳保手套 跨平台评论级扫描", font_size=32, bold=True)

# Subtitle info bar
add_textbox(slide2, 0.5, 1.35, 11, 0.3, "采集平台：抖音 / 小红书 / B站 / 微博 / 快手（失败）· 工具：OpenCLI + Yahoo Search · 2026-05-09", font_size=10, color=GRAY)

# Big numbers row
add_big_number(slide2, 0.5, 2.0, "10,000+", "抖音相关视频")
add_big_number(slide2, 3.2, 2.0, "5,000-20,000", "小红书相关笔记")
add_big_number(slide2, 6.0, 2.0, "92 万", "B站头部播放量")
add_big_number(slide2, 8.8, 2.0, "1,704 万", "微博超话阅读")

# Divider
add_rect(slide2, 0.5, 3.2, 12.3, 0.008, RGBColor(0xDD, 0xDD, 0xDD))

# Core finding
add_textbox(slide2, 0.5, 3.4, 12, 0.4, "核心结论", font_size=14, bold=True, color=CIOSH_GREEN)
add_textbox(slide2, 0.5, 3.75, 12, 0.6, "劳保手套的 C 端和 B 端需求都真实、活跃且体量可观。仅手套一个品类，就能激活庞大的交易需求。", font_size=16, bold=True)

# Key metrics
metrics = [
    ("评论采购信号占比", "抖音 40-60%  ·  小红书 30-50%（B2B 笔记达 100%）"),
    ("供应商竞价密度", "单条「找工厂」笔记 → 10 家供应商来自 5 省主动报价"),
    ("价格核心带", "C 端 ¥6.9-99  ·  B 端批发 ¥387-415（600双）"),
    ("头部互动量级", "抖音 3.5-11 万赞  ·  小红书 2.2-6.8 万赞  ·  B站 92 万播放"),
]
y = 4.55
for label, value in metrics:
    add_textbox(slide2, 0.7, y, 3.5, 0.35, label, font_size=11, bold=True, color=DARK)
    add_textbox(slide2, 4.2, y, 8, 0.35, value, font_size=11, color=DARK)
    y += 0.38

# Footer
add_textbox(slide2, 0.5, 7.0, 12, 0.3, "Messe Düsseldorf China  |  CIOSH Project  ·  S07", font_size=8, color=GRAY)

# ============================================================
# SLIDE 3: S07 平台维度详情
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
add_rect(slide3, 0, 0, 13.33, 0.06, CIOSH_GREEN)
add_rect(slide3, 0, 0.06, 0.06, 7.44, CIOSH_ORANGE)

add_textbox(slide3, 0.5, 0.25, 2, 0.3, "S07 · 平台扫描详情", font_size=10, bold=True, color=CIOSH_GREEN)
add_textbox(slide3, 0.5, 0.55, 11, 0.6, "五平台需求信号对比", font_size=26, bold=True)

# Platform cards - 2 rows x 3 cols layout
platforms = [
    ("抖音 ⭐⭐⭐⭐⭐", "B2B 直销主战场", [
        "头部视频 3.5-11万赞",
        "「厂家直销」「批发价」出镜率极高",
        "小黄车闭环成交",
        "品牌矩阵运营（华诚/星宇/护手匠等）",
        "使用场景：工地>汽修>工厂>搬运",
    ]),
    ("小红书 ⭐⭐⭐⭐⭐", "数据最扎实的 B2B 采购场", [
        "B2B 采购帖5赞→10家供应商5省竞价",
        "女性版型手套 = 明确市场空白",
        "Ansell 官方现身 1912 赞笔记",
        "「有链接吗」「透气吗」采购信号密集",
        "B2B清库存帖直接成交",
    ]),
    ("B站 ⭐⭐⭐⭐", "专业内容深度最好", [
        "头部视频 92.3万播放（工业安全专家）",
        "EN388 标准检测视频 + 留电话直销",
        "「大厂采购vs家族采购」B2B讨论",
        "UP主: 工业安全/测评/工厂号/健身跨界",
        "KOL 高某人-职业病防护 92万+播放",
    ]),
    ("微博 ⭐⭐⭐", "超话社区小但精准", [
        "手套超话 1,704万阅读 / 631帖",
        "「机工战略」79.7万粉机构号发声",
        "「了解行情,劳保路上不迷路」定位精准",
        "局限：帖量偏低，非B2B交易主战场",
    ]),
    ("快手", "采集失败（父agent超时）", [
        "下沉市场关键平台",
        "建议后续补充采集",
        "—",
        "—",
        "—",
    ]),
    ("知乎", "采集失败（需登录态）", [
        "OpenCLI 返回空数组",
        "裸浏览器触发反爬",
        "—",
        "—",
        "—",
    ]),
]

card_w = 3.85
card_h = 3.0
start_x = 0.5
start_y = 1.35
gap = 0.18

for i, (name, subtitle, points) in enumerate(platforms):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap)
    y = start_y + row * (card_h + gap)
    
    # Card background
    card = slide3.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(card_w), Inches(card_h)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = LIGHT_GRAY
    card.line.fill.background()
    
    # Platform name
    add_textbox(slide3, x+0.15, y+0.1, card_w-0.3, 0.3, name, font_size=12, bold=True, color=CIOSH_GREEN)
    add_textbox(slide3, x+0.15, y+0.38, card_w-0.3, 0.2, subtitle, font_size=9, color=GRAY)
    
    # Points
    py = y + 0.65
    for pt in points:
        add_textbox(slide3, x+0.15, py, card_w-0.3, 0.25, f"· {pt}", font_size=9, color=DARK)
        py += 0.25

add_textbox(slide3, 0.5, 7.0, 12, 0.3, "Messe Düsseldorf China  |  CIOSH Project  ·  S07", font_size=8, color=GRAY)

# ============================================================
# SLIDE 4: S07 小红书深度 + 战略启示
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
add_rect(slide4, 0, 0, 13.33, 0.06, CIOSH_GREEN)
add_rect(slide4, 0, 0.06, 0.06, 7.44, CIOSH_ORANGE)

add_textbox(slide4, 0.5, 0.25, 3, 0.3, "S07 · 小红书深度拆解", font_size=10, bold=True, color=CIOSH_GREEN)
add_textbox(slide4, 0.5, 0.55, 11, 0.6, "供应商竞价案例：一个5赞帖子引发的B2B争夺战", font_size=22, bold=True)

# Left: Case detail
add_textbox(slide4, 0.5, 1.3, 6, 0.3, "帖子「寻找劳保手套工厂」· 5赞 · 10条评论 全为供应商竞价", font_size=12, bold=True, color=CIOSH_ORANGE)

suppliers = [
    ("南通志柳劳护手套", "江苏", "「预算低找山东，要求高找南通」"),
    ("小红薯286555D3", "山东", "「我是山东的，你来找我」"),
    ("劳保手套工厂店", "浙江", "「源头直销，先看货再决定」"),
    ("诗晨", "山东", "「工厂自产自销，没有中间商」"),
    ("丁亿手套", "浙江", "「我这里有现货哦」"),
    ("恒顺劳保手套厂", "河北", "（参与竞价）"),
    ("Safety沃商防护用品", "广东", "（参与竞价）"),
]

for i, (name, loc, quote) in enumerate(suppliers):
    y = 1.75 + i * 0.35
    add_textbox(slide4, 0.7, y, 2.0, 0.3, name, font_size=10, bold=True, color=DARK)
    add_textbox(slide4, 2.8, y, 0.8, 0.3, loc, font_size=10, color=GRAY)
    add_textbox(slide4, 3.7, y, 3.0, 0.3, quote, font_size=10, color=DARK)

add_textbox(slide4, 0.5, 4.2, 6.5, 0.5, "结论：5赞帖子从 江苏/山东/浙江/河北/广东 五省吸引供应商主动报价。小红书上 B2B 采购需求匹配真实、供应商端高度活跃。", font_size=11, bold=True, color=CIOSH_GREEN)

# Right: 信息不对称痛点
add_rect(slide4, 7.3, 1.3, 5.5, 0.05, CIOSH_ORANGE)
add_textbox(slide4, 7.3, 1.4, 5.5, 0.35, "信息不对称痛点 = 展会价值点", font_size=13, bold=True, color=DARK)

pain_points = [
    ("「怎么买，淘宝搜不到」", "→ 展商信息无法穿透到买家"),
    ("「某多电焊手套全是假的」", "→ 供应链诚信危机，买家需可信渠道"),
    ("「预算低找山东，要求高找南通」", "→ 买家不知如何按品质/产地筛选供应商"),
    ("「均码只适合比较壮的手」", "→ 产品规格与真实需求脱节（女性手套空白）"),
    ("「这手套哪里买的」6.8万赞", "→ 爆款内容 vs 购买路径断裂"),
]
y = 1.85
for pain, implication in pain_points:
    add_textbox(slide4, 7.3, y, 5.8, 0.3, pain, font_size=10, bold=True, color=DARK)
    add_textbox(slide4, 7.3, y+0.22, 5.8, 0.25, implication, font_size=9, color=CIOSH_GREEN)
    y += 0.5

# Bottom: 产品机会
add_rect(slide4, 7.3, 4.35, 5.5, 0.05, CIOSH_ORANGE)
add_textbox(slide4, 7.3, 4.45, 5.5, 0.3, "产品机会信号", font_size=13, bold=True, color=DARK)

opps = [
    "女性版型劳保手套 — 4.6万赞笔记揭示市场空白",
    "「不闷手」透气手套 — 抖音 3.5万赞，夏季刚需",
    "平价替代进口 — Ansell 平替需求（1912赞）",
]
for i, opp in enumerate(opps):
    add_textbox(slide4, 7.3, 4.8 + i*0.3, 5.8, 0.25, f"· {opp}", font_size=10, color=DARK)

# Strategic insight
add_textbox(slide4, 0.5, 5.3, 12, 0.4, "战略启示：平台优先级", font_size=14, bold=True, color=DARK)

platform_pri = [
    ("抖音 ⭐⭐⭐⭐⭐", "工厂短视频 + 小黄车 + 批发线索"),
    ("小红书 ⭐⭐⭐⭐⭐", "B2B关键词布局 + 供应商竞价 + 女性品类"),
    ("B站 ⭐⭐⭐⭐", "工业安全科普 + EN388解读 + KOL合作"),
    ("1688 ⭐⭐⭐⭐", "展商线上店 + 搜索竞价"),
]

for i, (pf, strategy) in enumerate(platform_pri):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.2
    y = 5.7 + row * 0.35
    add_textbox(slide4, x, y, 2.5, 0.3, pf, font_size=11, bold=True, color=CIOSH_GREEN)
    add_textbox(slide4, x+2.6, y, 3.5, 0.3, strategy, font_size=10, color=DARK)

add_textbox(slide4, 0.5, 7.0, 12, 0.3, "Messe Düsseldorf China  |  CIOSH Project  ·  S07", font_size=8, color=GRAY)


# ============================================================
# SLIDE 5: S08 Sommelier.bot 竞品拆解
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
add_rect(slide5, 0, 0, 13.33, 0.06, CIOSH_GREEN)
add_rect(slide5, 0, 0.06, 0.06, 7.44, CIOSH_ORANGE)

add_textbox(slide5, 0.5, 0.25, 2, 0.3, "S08 · 竞品分析", font_size=10, bold=True, color=CIOSH_GREEN)
add_textbox(slide5, 0.5, 0.55, 11, 0.6, "Sommelier.bot 拆解 → CIOSH 功能映射", font_size=26, bold=True)

# Left: Product profile
add_rect(slide5, 0.5, 1.2, 5.8, 0.05, CIOSH_ORANGE)
add_textbox(slide5, 0.5, 1.3, 5.8, 0.3, "产品画像", font_size=14, bold=True, color=DARK)

profile = [
    "行业：红酒/烈酒/啤酒 B2B + B2C",
    "定位：24/7 AI 销售助手，嵌入商户网站",
    "规模：10万+终端用户 / 40+入驻商户 / 5个国家",
    "定价：免费 / €299/月 / Enterprise 定制",
    "部署：一行 JS 嵌入商户网站，一周上线",
    "核心叙事：「品类筛选已死，你需要 AI 侍酒师」",
]
for i, item in enumerate(profile):
    add_textbox(slide5, 0.7, 1.7 + i*0.3, 5.5, 0.25, f"· {item}", font_size=10, color=DARK)

# CIOSH translation
add_rect(slide5, 0.5, 3.6, 5.8, 0.05, CIOSH_GREEN)
add_textbox(slide5, 0.5, 3.7, 5.8, 0.4, "翻译成 CIOSH 语言", font_size=14, bold=True, color=DARK)
add_textbox(slide5, 0.5, 4.1, 5.8, 0.5, "「展商名录已死。展会官网没人看。你的观众/买家需要一个 AI 展会助手。」", font_size=13, bold=True, color=CIOSH_GREEN)

# Right: Architecture
add_rect(slide5, 7.0, 1.2, 5.8, 0.05, CIOSH_ORANGE)
add_textbox(slide5, 7.0, 1.3, 5.8, 0.3, "三层架构 → CIOSH 对标", font_size=14, bold=True, color=DARK)

layers = [
    ("B2C 免费平台", "app.sommelier.bot → CIOSH 公开 H5 Chat", "所有商户产品池，消费者搜索推荐，免费入驻"),
    ("商户独立实例 €299/月", "Merchant Instance → 展商专属 Agent", "单一商户专属，品牌定制，对话洞察，一行JS嵌入"),
    ("Enterprise 定制", "Enterprise → MDS 内部 + 大展商", "ERP/CRM对接，人工客服转接"),
]
y = 1.7
for name, mapping, desc in layers:
    add_textbox(slide5, 7.0, y, 2.8, 0.25, name, font_size=11, bold=True, color=DARK)
    add_textbox(slide5, 7.0, y+0.25, 5.8, 0.2, mapping, font_size=9, color=CIOSH_GREEN)
    add_textbox(slide5, 7.0, y+0.45, 5.8, 0.2, desc, font_size=9, color=GRAY)
    y += 0.75

# Bottom: Key insights
add_rect(slide5, 0.5, 4.9, 12.3, 0.05, CIOSH_GREEN)
add_textbox(slide5, 0.5, 5.0, 12, 0.3, "三点核心启示", font_size=14, bold=True, color=DARK)

insights = [
    ("启示 1", "B2C 免费是获客引擎，不是产品 — 免费层给商户带流量 → 商户看到价值 → 付费升级"),
    ("启示 2", "「一行 JS」是这个品类的标准做法 — B2B 展商最怕技术门槛，URL/JS Widget 是最低摩擦部署"),
    ("启示 3", "AI 富化数据是核心护城河 — 把烂数据变成好数据（30+属性标签）的价值远超 Chat 本身"),
]
for i, (title, desc) in enumerate(insights):
    y = 5.4 + i * 0.5
    add_textbox(slide5, 0.7, y, 1.0, 0.3, title, font_size=11, bold=True, color=CIOSH_ORANGE)
    add_textbox(slide5, 1.8, y, 10.5, 0.4, desc, font_size=10, color=DARK)

add_textbox(slide5, 0.5, 7.0, 12, 0.3, "Messe Düsseldorf China  |  CIOSH Project  ·  S08", font_size=8, color=GRAY)

# ============================================================
# SLIDE 6: S08 功能映射 + 差距 + 优先级
# ============================================================
slide6 = prs.slides.add_slide(blank_layout)
add_rect(slide6, 0, 0, 13.33, 0.06, CIOSH_GREEN)
add_rect(slide6, 0, 0.06, 0.06, 7.44, CIOSH_ORANGE)

add_textbox(slide6, 0.5, 0.25, 3, 0.3, "S08 · 功能映射 + 差距评估", font_size=10, bold=True, color=CIOSH_GREEN)
add_textbox(slide6, 0.5, 0.55, 11, 0.6, "CIOSH 功能对标 & 实施优先级", font_size=22, bold=True)

# Feature mapping table
add_textbox(slide6, 0.5, 1.2, 12, 0.3, "用户 / 展商 / 平台三侧功能映射", font_size=13, bold=True, color=CIOSH_GREEN)

mappings = [
    ("Sommelier.bot", "面向谁", "CIOSH 对等功能", "实现方式"),
    ("B2C 免费搜索", "消费者", "CIOSH 公开 H5 Chat", "H5 + Hermes Profile"),
    ("AI 产品推荐", "消费者", "「帮我找南通丁腈手套厂」", "展商数据 + LLM 匹配"),
    ("交互式葡萄酒地图", "消费者", "展馆平面图 + 展商标注", "H5 地图组件"),
    ("上传产品库存", "商户", "展商上传产品清单", "CSV上传 + 表单"),
    ("AI 富化 30+属性", "商户", "自动标签品类/产地/认证", "Skill 结构化处理"),
    ("品牌定制 Chat", "商户", "展商专属 Personality", "Profile + AGENTS.md"),
    ("JS 一行嵌入", "商户", "URL / JS Widget", "H5 Chat Widget"),
    ("对话洞察 Dashboard", "商户", "周报推送", "Cronjob 自动生成"),
    ("商户管理后台", "平台", "展商入驻管理", "Hermes 内部 Profile"),
]

y = 1.55
for i, row_data in enumerate(mappings):
    is_header = (i == 0)
    bg_color = CIOSH_GREEN if is_header else (LIGHT_GRAY if i % 2 == 0 else WHITE)
    font_color = WHITE if is_header else DARK
    font_bold = is_header
    
    col_widths = [2.6, 1.2, 3.0, 3.2]
    col_starts = [0.6, 3.2, 4.5, 7.5]
    
    for j, (text, w, x) in enumerate(zip(row_data, col_widths, col_starts)):
        add_textbox(slide6, x, y, w, 0.25, text, font_size=9 if not is_header else 10,
                    bold=font_bold, color=font_color)
    y += 0.25

# Right side: 差距评估
add_textbox(slide6, 0.5, 4.3, 6, 0.3, "差距评估", font_size=13, bold=True, color=CIOSH_GREEN)

gap_data = [
    ("公开 AI 对话", "✅ S.bot", "❌ CIOSH", "需搭建 H5 Chat"),
    ("AI 数据富化", "✅ 30+属性", "❌ 无", "需开发 Skill"),
    ("商户独立实例", "✅ €299/月", "❌ 无", "需展商 Profile 模板"),
    ("JS 嵌入", "✅ 一行script", "❌ 无", "需 Chat Widget"),
    ("多平台自动采集", "❌ 无", "✅ OpenCLI 扫描", "反超！"),
    ("主动信号捕获", "❌ 纯被动", "✅ S07 已验证", "反超！差异化武器"),
]

for i, (cap, sb, ciosh, gap) in enumerate(gap_data):
    y = 4.65 + i * 0.3
    add_textbox(slide6, 0.7, y, 2.0, 0.25, cap, font_size=9, bold=True, color=DARK)
    add_textbox(slide6, 2.8, y, 1.5, 0.25, sb, font_size=9, color=GRAY)
    add_textbox(slide6, 4.5, y, 1.5, 0.25, ciosh, font_size=9, color=GRAY)
    clr = CIOSH_GREEN if "反超" in gap else CIOSH_ORANGE
    add_textbox(slide6, 6.2, y, 3.5, 0.25, gap, font_size=9, color=clr)

# Priority section
add_textbox(slide6, 0.5, 6.1, 6, 0.3, "实施优先级", font_size=13, bold=True, color=CIOSH_GREEN)

priorities = [
    ("P0", "展商数据结构化 + AI 富化 · CIOSH 公开 H5 Chat"),
    ("P1", "Chat 可查询展商数据 · Widget 可嵌入（URL/JS）"),
    ("P2", "展商专属 Agent Profile 模板 · 自媒体信号捕获匹配"),
    ("P3", "对话分析周报 · 品牌定制（颜色/语调）"),
]

for i, (level, task) in enumerate(priorities):
    y = 6.4 + i * 0.28
    add_textbox(slide6, 0.7, y, 0.5, 0.25, level, font_size=10, bold=True, color=CIOSH_ORANGE)
    add_textbox(slide6, 1.3, y, 10, 0.25, task, font_size=9, color=DARK)

add_textbox(slide6, 0.5, 7.0, 12, 0.3, "Messe Düsseldorf China  |  CIOSH Project  ·  S08", font_size=8, color=GRAY)

# ============================================================
# Save
# ============================================================
output_path = "/Volumes/databoard/AI Project/ciosh/output/CIOSH_Blackbox_拆解路径.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")

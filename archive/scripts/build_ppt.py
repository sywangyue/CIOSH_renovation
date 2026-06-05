#!/usr/bin/env python3
"""Build CIOSH Blackbox one-page PPT slide on MDS brand template."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ── Constants ──────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# MDS Brand Colors
DARK_BLUE   = RGBColor(0x0E, 0x28, 0x41)
MED_BLUE    = RGBColor(0x15, 0x60, 0x82)
ORANGE      = RGBColor(0xE9, 0x71, 0x32)
LIGHT_GRAY  = RGBColor(0xE8, 0xE8, 0xE8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x99, 0x99, 0x99)
AMBER       = RGBColor(0xF0, 0xAD, 0x4E)
GREEN       = RGBColor(0x19, 0x6B, 0x24)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
LIGHT_AMBER = RGBColor(0xFD, 0xF0, 0xD5)
LIGHT_GREEN = RGBColor(0xE2, 0xF0, 0xD9)
CARD_BG     = RGBColor(0xF5, 0xF7, 0xFA)
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)

# Font
FONT_CN = '微软雅黑'
FONT_EN = 'Arial'

# ── Helpers ────────────────────────────────────────────────────

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(10),
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=None, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    """Add a text box and return (shape, text_frame)."""
    if font_name is None:
        font_name = FONT_CN
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # Set anchor
    txBox.text_frame.paragraphs[0].alignment = alignment
    # Set vertical anchor via XML
    bodyPr = txBox.text_frame._txBody.find(qn('a:bodyPr'))
    if bodyPr is not None:
        anchor_map = {
            MSO_ANCHOR.TOP: 't',
            MSO_ANCHOR.MIDDLE: 'ctr',
            MSO_ANCHOR.BOTTOM: 'b',
        }
        bodyPr.set('anchor', anchor_map.get(anchor, 't'))

    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    # Set East Asian font
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), font_name if font_name else FONT_CN)
    return txBox, tf

def add_paragraph(tf, text, font_size=Pt(10), font_color=BLACK, bold=False,
                  alignment=PP_ALIGN.LEFT, font_name=None, space_before=Pt(2),
                  space_after=Pt(2)):
    """Add a paragraph to existing text frame."""
    if font_name is None:
        font_name = FONT_CN
    p = tf.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), font_name if font_name else FONT_CN)
    return p

def add_rounded_rect(slide, left, top, width, height, fill_color=DARK_BLUE,
                     border_color=None, border_width=Pt(0)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color and border_width > 0:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.1
    return shape

def add_arrow_right(slide, left, top, width, height, color=MED_BLUE):
    """Add a right-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height, color=MED_BLUE):
    """Add a down-pointing arrow."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_chevron(slide, left, top, width, height, color=ORANGE):
    """Add a chevron shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def set_shape_text(shape, text, font_size=Pt(8), font_color=WHITE,
                   bold=False, alignment=PP_ALIGN.CENTER, font_name=None):
    """Set text on an auto shape."""
    if font_name is None:
        font_name = FONT_CN
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    rPr = run._r.get_or_add_rPr()
    rPr.set(qn('a:ea'), font_name if font_name else FONT_CN)

# ── Main Build ─────────────────────────────────────────────────

def build_slide():
    prs = Presentation('/Volumes/Maxgo/Branding/PPT Template_MDS_0214.pptx')

    # Use blank layout (index 7)
    blank_layout = prs.slide_layouts[7]
    slide = prs.slides.add_slide(blank_layout)

    # Remove any existing placeholder shapes from the blank layout
    # (they're typically image placeholders)
    shapes_to_remove = []
    for shape in slide.shapes:
        if shape.is_placeholder:
            shapes_to_remove.append(shape)
    for s in shapes_to_remove:
        sp = s._element
        sp.getparent().remove(sp)

    # ═══════════════════════════════════════════════════════════
    # LAYER 0: Full background — white
    # ═══════════════════════════════════════════════════════════
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # ═══════════════════════════════════════════════════════════
    # LAYER 1: Title Bar (0.00 – 0.75")
    # ═══════════════════════════════════════════════════════════
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.75)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()

    # MDS logo area (left side accent)
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.08), Inches(0.75)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ORANGE
    accent_bar.line.fill.background()

    # Title text
    add_textbox(slide, Inches(0.4), Inches(0.08), Inches(9), Inches(0.42),
                "CIOSH 痛点诊断与 Blackbox 拆解路径",
                font_size=Pt(28), font_color=WHITE, bold=True,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle
    add_textbox(slide, Inches(0.4), Inches(0.42), Inches(9), Inches(0.28),
                "从调研诊断到假设验证 — 信息黑盒的识别、描述与应对框架",
                font_size=Pt(11), font_color=RGBColor(0xBB, 0xCC, 0xDD), bold=False,
                alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

    # Date tag (右上角)
    add_textbox(slide, Inches(10.5), Inches(0.48), Inches(2.5), Inches(0.22),
                "2026.05.09  |  CIOSH Project",
                font_size=Pt(8), font_color=RGBColor(0x99, 0xAA, 0xBB), bold=False,
                alignment=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    # ═══════════════════════════════════════════════════════════
    # LAYER 2: Research Timeline (0.85" – 1.70")
    # ═══════════════════════════════════════════════════════════
    timeline_y = Inches(0.92)

    # Section label
    add_textbox(slide, Inches(0.4), timeline_y, Inches(1.6), Inches(0.24),
                "调研路径", font_size=Pt(9), font_color=MED_GRAY, bold=True,
                alignment=PP_ALIGN.LEFT)

    nodes = [
        ("S01", "品类断层图谱\n4/24"),
        ("S03", "POC 变革方案\n4/28"),
        ("S04", "劳保云调研\n4/30"),
        ("S05", "内部访谈\n5/7"),
        ("S06", "小程序聚合\n5/8"),
    ]

    node_w = Inches(2.0)
    node_h = Inches(0.55)
    arrow_w = Inches(0.25)
    arrow_h = Inches(0.16)
    start_x = Inches(0.4)
    gap = Inches(0.12)
    node_y = Inches(1.2)

    for i, (code, label) in enumerate(nodes):
        x = start_x + i * (node_w + arrow_w + gap)

        # Node
        is_last = (i == len(nodes) - 1)
        node_color = ORANGE if is_last else DARK_BLUE
        rect = add_rounded_rect(slide, x, node_y, node_w, node_h,
                                fill_color=node_color, border_color=None)
        # Code
        _, tf = add_textbox(slide, x + Inches(0.08), node_y + Inches(0.04),
                            Inches(0.5), Inches(0.22),
                            code, font_size=Pt(9), font_color=WHITE, bold=True,
                            alignment=PP_ALIGN.LEFT)
        # Label
        add_textbox(slide, x + Inches(0.6), node_y + Inches(0.04),
                    node_w - Inches(0.7), node_h - Inches(0.08),
                    label, font_size=Pt(7.5), font_color=WHITE, bold=False,
                    alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

        # Arrow between nodes
        if not is_last:
            ax = x + node_w + Inches(0.02)
            add_arrow_right(slide, ax, node_y + Inches(0.2), arrow_w, arrow_h,
                            color=MED_BLUE)

    # ═══════════════════════════════════════════════════════════
    # LAYER 3: Three Interview Ports (1.85" – 3.10")
    # ═══════════════════════════════════════════════════════════
    ports_y = Inches(2.0)
    add_textbox(slide, Inches(0.4), ports_y, Inches(3), Inches(0.22),
                "三端口调研 (5/7–5/8 访谈输出)", font_size=Pt(9),
                font_color=MED_GRAY, bold=True)

    col_w = Inches(3.8)
    col_h = Inches(0.95)
    col_gap = Inches(0.25)
    col_start_x = Inches(0.4)
    col_y = Inches(2.28)
    port_colors = [MED_BLUE, DARK_BLUE, ORANGE]
    port_letters = ['A', 'B', 'C']
    port_titles = [
        "人口相传",
        "全体大会",
        "端口调研"
    ]
    port_subs = [
        "Rya / James / Kate 一手信息",
        "中方参与 · 更多声音",
        "品类缺口 + 技术规范"
    ]
    port_details = [
        "项目运行机制 SOP\n数据敏感性盘点\n双方分工与权限结构",
        "痛点结构化(六条根因)\nPOC 变革方案(六个)\n信息入口→流通→现场三层模型",
        "8×4 品类断层矩阵\n劳保云/SaaS-B 技术栈\n三端集成断链分析"
    ]

    for i in range(3):
        x = col_start_x + i * (col_w + col_gap)

        # Card background
        card = add_rounded_rect(slide, x, col_y, col_w, col_h,
                                fill_color=CARD_BG, border_color=None)

        # Left accent bar on card
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, col_y, Inches(0.06), col_h
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = port_colors[i]
        accent.line.fill.background()

        # Port letter badge
        badge_size = Inches(0.28)
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x + Inches(0.16), col_y + Inches(0.08),
            badge_size, badge_size
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = port_colors[i]
        badge.line.fill.background()
        set_shape_text(badge, port_letters[i], font_size=Pt(11), font_color=WHITE,
                       bold=True, font_name=FONT_EN)

        # Title
        add_textbox(slide, x + Inches(0.52), col_y + Inches(0.06),
                    Inches(2.0), Inches(0.22),
                    port_titles[i], font_size=Pt(12), font_color=port_colors[i],
                    bold=True)
        # Subtitle
        add_textbox(slide, x + Inches(0.52), col_y + Inches(0.27),
                    Inches(3.1), Inches(0.18),
                    port_subs[i], font_size=Pt(7), font_color=MED_GRAY)

        # Detail text
        add_textbox(slide, x + Inches(0.18), col_y + Inches(0.48),
                    col_w - Inches(0.36), Inches(0.44),
                    port_details[i], font_size=Pt(7.5), font_color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT, line_spacing=1.2)

    # Down arrows from each column → center
    arrow_down_y = Inches(3.26)
    for i in range(3):
        ax = col_start_x + i * (col_w + col_gap) + col_w / 2 - Inches(0.10)
        add_down_arrow(slide, ax, arrow_down_y, Inches(0.20), Inches(0.22),
                       color=MED_GRAY)

    # ═══════════════════════════════════════════════════════════
    # LAYER 4: BLACKBOX Central (3.50" – 4.85")
    # ═══════════════════════════════════════════════════════════
    bb_y = Inches(3.58)
    bb_w = Inches(8.0)
    bb_h = Inches(1.15)
    bb_x = (Inches(13.333) - bb_w) / 2

    # Blackbox background
    bb = add_rounded_rect(slide, bb_x, bb_y, bb_w, bb_h,
                          fill_color=DARK_BLUE,
                          border_color=ORANGE, border_width=Pt(2.5))

    # "BLACKBOX" label
    add_textbox(slide, bb_x, bb_y + Inches(0.06), bb_w, Inches(0.30),
                "B L A C K B O X",
                font_size=Pt(20), font_color=ORANGE, bold=True,
                alignment=PP_ALIGN.CENTER, font_name=FONT_EN,
                anchor=MSO_ANCHOR.MIDDLE)

    # Divider line
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        bb_x + Inches(2.5), bb_y + Inches(0.42),
        bb_w - Inches(5.0), Inches(0.015)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = ORANGE
    div.line.fill.background()

    # Description lines
    desc_lines = [
        ("核心问题: ", "展商 ↔ 观众 ↔ 主办方  三端信息完全断裂，看不到内部流转机制"),
        ("关键特征: ", "可以描述 Blackbox 的输入/输出特征，但无法穿透其内部"),
        ("小程序定位: ", "小程序可能是解决黑盒的一种方式，但不是全部"),
    ]
    for j, (label, desc) in enumerate(desc_lines):
        line_y = bb_y + Inches(0.52) + j * Inches(0.20)
        txBox, tf = add_textbox(slide, bb_x + Inches(0.8), line_y,
                                bb_w - Inches(1.6), Inches(0.18),
                                "", font_size=Pt(9), font_color=WHITE,
                                anchor=MSO_ANCHOR.MIDDLE)
        p = tf.paragraphs[0]
        run_label = p.add_run()
        run_label.text = label
        run_label.font.size = Pt(9)
        run_label.font.color.rgb = ORANGE
        run_label.font.bold = True
        run_label.font.name = FONT_CN
        rPr = run_label._r.get_or_add_rPr()
        rPr.set(qn('a:ea'), FONT_CN)

        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(8.5)
        run_desc.font.color.rgb = WHITE
        run_desc.font.bold = False
        run_desc.font.name = FONT_CN
        rPr2 = run_desc._r.get_or_add_rPr()
        rPr2.set(qn('a:ea'), FONT_CN)

    # ═══════════════════════════════════════════════════════════
    # LAYER 5: Three Hypotheses (4.95" – 6.70")
    # ═══════════════════════════════════════════════════════════
    hyp_y_label = Inches(5.0)
    add_textbox(slide, Inches(0.4), hyp_y_label, Inches(4), Inches(0.22),
                "三种假设路径 — 对 Blackbox 的应对策略",
                font_size=Pt(9), font_color=MED_GRAY, bold=True)

    hyp_w = Inches(3.95)
    hyp_h = Inches(1.5)
    hyp_gap = Inches(0.2)
    hyp_start_x = Inches(0.4)
    hyp_y = Inches(5.25)

    hypotheses = [
        {
            'num': '路径 1',
            'title': '穿透黑盒',
            'prob': '概率 ≈ 0%',
            'method': '直接打通现有系统\n劳保云 / SaaS-B / 官网',
            'risk': '组织壁垒不可逾越\n数据主权分散，技术接口不存在',
            'verdict': '❌ 不可行',
            'fill': CARD_BG,
            'accent': MED_GRAY,
            'prob_color': RED_ACCENT,
            'verdict_color': MED_GRAY,
        },
        {
            'num': '路径 2',
            'title': '绕开黑盒',
            'prob': '概率 = 100%',
            'method': '绕过所有现有系统\n独立建设新信息通道',
            'risk': '信息产生波动和发散\n衍生更多微型黑盒，长期复杂度上升',
            'verdict': '⚠ 可行但发散',
            'fill': LIGHT_AMBER,
            'accent': AMBER,
            'prob_color': AMBER,
            'verdict_color': AMBER,
        },
        {
            'num': '路径 3',
            'title': '白名单冲抵',
            'prob': '可行度 高 ★',
            'method': '自建 Whitelist 信息流\n冲抵 Blackbox 后 merge 回整体',
            'risk': '工程和技术难度高\n需要精巧的信息架构设计',
            'verdict': '✓ 推荐方向',
            'fill': LIGHT_GREEN,
            'accent': GREEN,
            'prob_color': GREEN,
            'verdict_color': GREEN,
        },
    ]

    for i, hyp in enumerate(hypotheses):
        x = hyp_start_x + i * (hyp_w + hyp_gap)

        # Card
        card = add_rounded_rect(slide, x, hyp_y, hyp_w, hyp_h,
                                fill_color=hyp['fill'],
                                border_color=hyp['accent'], border_width=Pt(1.5))

        # Top accent strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, hyp_y, hyp_w, Inches(0.04)
        )
        strip.fill.solid()
        strip.fill.fore_color.rgb = hyp['accent']
        strip.line.fill.background()

        # Path number + Title
        add_textbox(slide, x + Inches(0.18), hyp_y + Inches(0.10),
                    hyp_w - Inches(0.36), Inches(0.24),
                    f"{hyp['num']}  {hyp['title']}",
                    font_size=Pt(14), font_color=DARK_BLUE, bold=True)

        # Probability badge
        prob_badge = add_rounded_rect(slide, x + Inches(0.18), hyp_y + Inches(0.38),
                                      Inches(1.8), Inches(0.24),
                                      fill_color=hyp['accent'])
        set_shape_text(prob_badge, hyp['prob'], font_size=Pt(10), font_color=WHITE,
                       bold=True)

        # Method
        add_textbox(slide, x + Inches(0.18), hyp_y + Inches(0.70),
                    hyp_w - Inches(0.36), Inches(0.32),
                    hyp['method'], font_size=Pt(8), font_color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT)

        # Risk label
        add_textbox(slide, x + Inches(0.18), hyp_y + Inches(1.02),
                    Inches(0.5), Inches(0.16),
                    "代价:", font_size=Pt(7), font_color=MED_GRAY, bold=True)

        # Risk text
        add_textbox(slide, x + Inches(0.55), hyp_y + Inches(1.02),
                    hyp_w - Inches(0.75), Inches(0.30),
                    hyp['risk'], font_size=Pt(7.5), font_color=DARK_GRAY,
                    alignment=PP_ALIGN.LEFT)

        # Verdict
        add_textbox(slide, x + hyp_w - Inches(1.5), hyp_y + Inches(0.10),
                    Inches(1.3), Inches(0.24),
                    hyp['verdict'], font_size=Pt(9), font_color=hyp['verdict_color'],
                    bold=True, alignment=PP_ALIGN.RIGHT)

    # ═══════════════════════════════════════════════════════════
    # LAYER 6: Bottom conclusion bar (7.00" – 7.50")
    # ═══════════════════════════════════════════════════════════
    conclusion_y = Inches(7.0)

    # Thin accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), conclusion_y, Inches(12.5), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ORANGE
    line.line.fill.background()

    # Conclusion text
    add_textbox(slide, Inches(0.4), conclusion_y + Inches(0.10),
                Inches(12.5), Inches(0.35),
                "结论: 路径3（白名单冲抵 → merge）为推荐方向。小程序是载体，Whitelist 是机制，merge 回完整信息流是目标。\n工程难度高但可行度高 — 需要在下一阶段完成信息架构设计 + 技术选型。",
                font_size=Pt(8), font_color=DARK_GRAY, bold=False,
                alignment=PP_ALIGN.LEFT)

    # Footer branding
    add_textbox(slide, Inches(10.0), conclusion_y + Inches(0.38),
                Inches(3.0), Inches(0.15),
                "Messe Düsseldorf China  |  CIOSH Project",
                font_size=Pt(6.5), font_color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT)

    # ── Remove template slides, keep only our new slide ──
    # Our slide is the last one; remove all others from the slide list
    pres_elem = prs._element
    sldIdLst = pres_elem.find(qn('p:sldIdLst'))
    sld_ids = list(sldIdLst)
    our_sld_id = sld_ids[-1]  # keep the last one
    for sld_id in sld_ids[:-1]:
        sldIdLst.remove(sld_id)

    # ── Save ──
    output_path = '/Volumes/databoard/AI Project/ciosh/output/CIOSH_Blackbox_拆解路径.pptx'
    prs.save(output_path)
    print(f"Saved to: {output_path}")
    return output_path

if __name__ == '__main__':
    build_slide()
